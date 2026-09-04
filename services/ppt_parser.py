from pptx import Presentation


class PPTParser:

    def extract(self, file_path: str):
        presentation = Presentation(file_path)

        slides = []

        for index, slide in enumerate(
                presentation.slides,
                start=1,
        ):
            title = ""

            if slide.shapes.title:
                title = slide.shapes.title.text.strip()

            content = []

            for shape in slide.shapes:
                if (
                        hasattr(shape, "text")
                        and shape.text.strip()
                        and shape != slide.shapes.title
                ):
                    content.append(
                        shape.text.strip()
                    )

            slides.append({
                "index": index,
                "title": title,
                "content": content,
            })

        return slides