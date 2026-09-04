from pptx import Presentation


class PPTModifier:

    def modify(
            self,
            file_path: str,
            modifications: list[dict],
            output_path: str,
    ):
        presentation = Presentation(file_path)

        for modification in modifications:
            action = modification.get("action")
            slide_index = modification.get("slide")

            if (
                    slide_index is not None
                    and (
                    slide_index < 1
                    or slide_index > len(presentation.slides)
            )
            ):
                continue

            slide = (
                presentation.slides[slide_index - 1]
                if slide_index
                else None
            )

            if action == "change_title" and slide:
                new_title = modification.get("value")

                if slide.shapes.title:
                    slide.shapes.title.text = new_title

            elif action == "replace_text":
                old_text = modification.get("old_text")
                new_text = modification.get("new_text")

                for current_slide in presentation.slides:
                    for shape in current_slide.shapes:
                        if not hasattr(shape, "text_frame"):
                            continue

                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if old_text in run.text:
                                    run.text = run.text.replace(
                                        old_text,
                                        new_text,
                                    )

        presentation.save(output_path)