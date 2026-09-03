from pathlib import Path

from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


def extract_pdf(file_path: str) -> list[dict]:
    reader = PdfReader(file_path)

    documents = []

    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""

        if text.strip():
            documents.append({
                "text": text.strip(),
                "source": f"Page {index}",
                "type": "page",
            })

    return documents


def extract_docx(file_path: str) -> list[dict]:
    document = Document(file_path)

    paragraphs = [
        paragraph.text.strip()
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]

    if not paragraphs:
        return []

    return [{
        "text": "\n".join(paragraphs),
        "source": "Document",
        "type": "document",
    }]


def extract_pptx(file_path: str) -> list[dict]:
    presentation = Presentation(file_path)

    documents = []

    for index, slide in enumerate(
            presentation.slides,
            start=1,
    ):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        if slide_text:
            documents.append({
                "text": "\n".join(slide_text),
                "source": f"Slide {index}",
                "type": "slide",
            })

    return documents


def extract_xlsx(file_path: str) -> list[dict]:
    workbook = load_workbook(
        file_path,
        read_only=True,
        data_only=True,
    )

    documents = []

    for worksheet in workbook.worksheets:
        rows = []

        for row in worksheet.iter_rows(
                values_only=True
        ):
            values = [
                str(value).strip()
                if value is not None
                else ""
                for value in row
            ]

            if any(values):
                rows.append(" | ".join(values))

        if rows:
            documents.append({
                "text": "\n".join(rows),
                "source": f"Sheet: {worksheet.title}",
                "type": "sheet",
            })

    return documents


def extract_document(file_path: str) -> list[dict]:
    extension = Path(file_path).suffix.lower()

    if extension == ".pdf":
        return extract_pdf(file_path)

    if extension == ".docx":
        return extract_docx(file_path)

    if extension == ".pptx":
        return extract_pptx(file_path)

    if extension in [".xlsx", ".xlsm"]:
        return extract_xlsx(file_path)

    raise ValueError(
        f"Unsupported file type: {extension}"
    )


def extract_text(file_path: str) -> str:
    documents = extract_document(file_path)

    return "\n\n".join(
        f"{item['source']}:\n{item['text']}"
        for item in documents
    )